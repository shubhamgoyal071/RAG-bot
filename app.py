"""
StudyBot — RAG-powered study assistant.
Supports PDF, DOCX, PPTX, TXT.
Supports Gemini and Groq as LLM providers.
"""

import io
import os
import pathlib
import tempfile

import chromadb
import streamlit as st
from chromadb.utils import embedding_functions
from dotenv import load_dotenv
from pypdf import PdfReader

try:
    import ppt2txt
    HAS_PPT2TXT = True
except ImportError:
    HAS_PPT2TXT = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# ── Config ─────────────────────────────────────────────────────────────────────
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
GROQ_API_KEY   = os.getenv("GROQ_API_KEY", "")

CHUNK_SIZE         = 1000
OVERLAP            = 150
DISTANCE_THRESHOLD = 1.1
PDF_DIR            = pathlib.Path("pdfs")
PDF_DIR.mkdir(exist_ok=True)

GEMINI_MODELS = {
    "Gemini 2.5 Flash": "gemini-2.5-flash",
    "Gemini 1.5 Flash": "gemini-1.5-flash",
    "Gemini 1.5 Pro":   "gemini-1.5-pro",
}
GROQ_MODELS = {
    "Llama 3.3 70B (recommended)": "llama-3.3-70b-versatile",
    "Llama 3.1 8B (fastest)":      "llama-3.1-8b-instant",
    "Mixtral 8x7B":                 "mixtral-8x7b-32768",
    "Gemma 2 9B":                   "gemma2-9b-it",
}

# ── Page setup — NO sidebar ────────────────────────────────────────────────────
st.set_page_config(
    page_title="StudyBot",
    page_icon="📖",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# ── CSS ────────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');

html, body, [class*="css"] {
  font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif !important;
  background: #09090b !important;
  color: #fafafa !important;
}

/* Hide ALL Streamlit chrome including sidebar toggle */
#MainMenu, footer, header, .stDeployButton,
[data-testid="stSidebar"],
[data-testid="collapsedControl"] { display: none !important; }

/* Remove default page padding */
.block-container {
  padding: 0 !important;
  max-width: 100% !important;
}

/* ── Column border — left panel ── */
[data-testid="stHorizontalBlock"] > div:first-child {
  border-right: 1px solid #1f1f23;
  background: #0c0c0e;
  min-height: calc(100vh - 56px);
  padding: 1.25rem 1rem !important;
}

/* ── Right panel ── */
[data-testid="stHorizontalBlock"] > div:last-child {
  padding: 1.25rem 1.5rem !important;
  background: #09090b;
}

/* ── Top nav ── */
.top-nav {
  display: flex; align-items: center; justify-content: space-between;
  padding: 0 1.5rem; height: 56px;
  border-bottom: 1px solid #1f1f23;
  background: #09090b;
}
.nav-brand { display: flex; align-items: center; gap: 9px; }
.nav-logo {
  width: 26px; height: 26px; background: #6366f1; border-radius: 6px;
  display: flex; align-items: center; justify-content: center; font-size: 13px;
}
.nav-name { font-size: 0.9rem; font-weight: 600; color: #fafafa; letter-spacing: -0.02em; }
.nav-tag {
  font-size: 0.61rem; font-weight: 600; color: #71717a;
  background: #18181b; border: 1px solid #27272a;
  border-radius: 4px; padding: 2px 6px; text-transform: uppercase;
}
.status-pill {
  display: inline-flex; align-items: center; gap: 6px;
  font-size: 0.75rem; color: #71717a;
}
.s-dot { width: 6px; height: 6px; border-radius: 50%; }

/* ── Section labels ── */
.sec-label {
  font-size: 0.67rem; font-weight: 600; color: #52525b;
  text-transform: uppercase; letter-spacing: 0.08em;
  padding-bottom: 0.5rem;
  display: block;
}

/* ── File uploader ── */
[data-testid="stFileUploader"] section {
  background: #18181b !important;
  border: 1.5px dashed #3f3f46 !important;
  border-radius: 10px !important;
  padding: 1rem 0.75rem !important;
  transition: border-color 0.2s !important;
  text-align: center !important;
}
[data-testid="stFileUploader"] section:hover {
  border-color: #6366f1 !important;
  background: rgba(99,102,241,0.04) !important;
}
[data-testid="stFileUploader"] section p,
[data-testid="stFileUploader"] section small {
  color: #71717a !important; font-size: 0.76rem !important;
}
[data-testid="stFileUploader"] section > div > div > span {
  color: #a1a1aa !important; font-size: 0.8rem !important; font-weight: 500 !important;
}
[data-testid="stFileUploader"] label { display: none !important; }

/* ── Format badges ── */
.fmt-row { display: flex; gap: 3px; margin-bottom: 0.6rem; flex-wrap: wrap; }
.fmt-badge {
  font-size: 0.62rem; font-weight: 700; letter-spacing: 0.04em;
  background: rgba(99,102,241,0.1); border: 1px solid rgba(99,102,241,0.22);
  color: #818cf8; border-radius: 4px; padding: 2px 6px; text-transform: uppercase;
}

/* ── Doc rows ── */
.doc-item {
  display: flex; align-items: center; gap: 7px;
  padding: 0.38rem 0.4rem; border-radius: 6px;
  transition: background 0.12s; margin-bottom: 1px;
}
.doc-item:hover { background: rgba(255,255,255,0.04); }
.doc-ico {
  width: 24px; height: 24px; min-width: 24px; border-radius: 5px;
  display: flex; align-items: center; justify-content: center; font-size: 11px;
}
.t-pdf  { background:rgba(239,68,68,.12);  border:1px solid rgba(239,68,68,.2);  }
.t-docx { background:rgba(59,130,246,.12); border:1px solid rgba(59,130,246,.2); }
.t-pptx { background:rgba(234,88,12,.12);  border:1px solid rgba(234,88,12,.2);  }
.t-txt  { background:#18181b; border:1px solid #27272a; }
.doc-name {
  font-size: 0.73rem; color: #a1a1aa;
  overflow: hidden; text-overflow: ellipsis; white-space: nowrap; flex: 1;
}
.doc-ok { width: 5px; height: 5px; border-radius: 50%; background: #22c55e; flex-shrink: 0; }

/* ── Stats ── */
.stats-grid { display: flex; gap: 8px; margin-top: 0.6rem; }
.stat-card {
  flex: 1; background: #18181b; border: 1px solid #27272a;
  border-radius: 8px; padding: 0.6rem 0.5rem; text-align: center;
}
.stat-n { font-size: 1.05rem; font-weight: 600; color: #fafafa; }
.stat-l { font-size: 0.62rem; color: #52525b; text-transform: uppercase; letter-spacing: 0.05em; margin-top: 1px; }

/* ── Selectbox ── */
.stSelectbox > div > div {
  background: #18181b !important;
  border: 1px solid #27272a !important;
  border-radius: 8px !important;
  color: #d4d4d8 !important;
  font-size: 0.83rem !important;
}
.stSelectbox label { color: #71717a !important; font-size: 0.74rem !important; font-weight: 500 !important; }

/* ── Provider badge ── */
.prov-tag {
  display: inline-flex; align-items: center; gap: 5px;
  font-size: 0.7rem; font-weight: 500; border-radius: 5px;
  padding: 3px 8px; margin-top: 4px; margin-bottom: 4px;
}
.p-groq   { background:rgba(249,115,22,.1); border:1px solid rgba(249,115,22,.25); color:#fb923c; }
.p-gemini { background:rgba(59,130,246,.1); border:1px solid rgba(59,130,246,.25); color:#60a5fa; }

/* ── Banners ── */
.banner {
  display: flex; align-items: flex-start; gap: 6px;
  padding: 0.5rem 0.75rem; border-radius: 7px;
  font-size: 0.76rem; line-height: 1.5; margin-bottom: 0.5rem;
}
.bw { background:rgba(234,179,8,.08);  border:1px solid rgba(234,179,8,.22);  color:#ca8a04; }
.bi { background:rgba(99,102,241,.08); border:1px solid rgba(99,102,241,.2);  color:#818cf8; }
.bo { background:rgba(34,197,94,.08);  border:1px solid rgba(34,197,94,.2);   color:#16a34a; }
.be { background:rgba(239,68,68,.08);  border:1px solid rgba(239,68,68,.22);  color:#f87171; }

/* ── Sidebar button ── */
[data-testid="stHorizontalBlock"] > div:first-child .stButton > button {
  background: transparent !important;
  border: 1px solid #27272a !important;
  color: #71717a !important;
  font-size: 0.79rem !important;
  font-weight: 500 !important;
  border-radius: 7px !important;
  width: 100% !important;
  padding: 0.45rem 0.8rem !important;
}
[data-testid="stHorizontalBlock"] > div:first-child .stButton > button:hover {
  background: #18181b !important; color: #d4d4d8 !important;
}

/* ── Tabs ── */
.stTabs [data-baseweb="tab-list"] {
  background: transparent !important; border-radius: 0 !important;
  border-bottom: 1px solid #27272a !important;
  padding: 0 !important; gap: 0 !important; margin-bottom: 0.5rem !important;
}
.stTabs [data-baseweb="tab"] {
  border-radius: 0 !important; color: #71717a !important;
  font-weight: 500 !important; font-size: 0.84rem !important;
  padding: 10px 16px !important; margin-bottom: -1px !important;
  border-bottom: 2px solid transparent !important;
  background: transparent !important;
  border-top: none !important; border-left: none !important; border-right: none !important;
  transition: color 0.15s !important;
}
.stTabs [data-baseweb="tab"]:hover  { color: #d4d4d8 !important; }
.stTabs [aria-selected="true"] {
  color: #fafafa !important; border-bottom-color: #6366f1 !important;
  background: transparent !important; box-shadow: none !important;
}
.stTabs [data-baseweb="tab-panel"] { padding: 0.75rem 0 0 !important; }

/* ── Inputs (main area) ── */
[data-testid="stHorizontalBlock"] > div:last-child .stTextInput > div > div > input,
[data-testid="stHorizontalBlock"] > div:last-child .stTextArea textarea {
  background: #18181b !important; border: 1px solid #27272a !important;
  border-radius: 9px !important; color: #fafafa !important;
  font-family: 'Inter', sans-serif !important; font-size: 0.875rem !important;
  caret-color: #6366f1 !important;
  transition: border-color 0.15s, box-shadow 0.15s !important;
}
[data-testid="stHorizontalBlock"] > div:last-child .stTextInput > div > div > input:focus,
[data-testid="stHorizontalBlock"] > div:last-child .stTextArea textarea:focus {
  border-color: #6366f1 !important;
  box-shadow: 0 0 0 3px rgba(99,102,241,0.1) !important;
}
[data-testid="stHorizontalBlock"] > div:last-child .stTextInput > div > div > input::placeholder { color: #52525b !important; }
[data-testid="stHorizontalBlock"] > div:last-child .stTextInput label,
[data-testid="stHorizontalBlock"] > div:last-child .stTextArea label { color: #71717a !important; font-size: 0.8rem !important; }

/* ── Main Send button ── */
[data-testid="stHorizontalBlock"] > div:last-child .stButton > button {
  background: #6366f1 !important; color: #fff !important; border: none !important;
  border-radius: 8px !important; font-family: 'Inter', sans-serif !important;
  font-weight: 500 !important; font-size: 0.855rem !important;
  padding: 0.58rem 1rem !important; width: 100% !important;
  transition: background 0.15s !important; white-space: nowrap !important;
}
[data-testid="stHorizontalBlock"] > div:last-child .stButton > button:hover {
  background: #4f46e5 !important;
}

/* ── Slider ── */
.stSlider [data-testid="stSliderThumb"] { background: #6366f1 !important; border-color: #6366f1 !important; }
.stSlider label { color: #71717a !important; font-size: 0.8rem !important; }

/* ── Chat messages ── */
.msg { display: flex; gap: 10px; margin-bottom: 0.9rem; animation: mi 0.18s ease; }
.msg.u { flex-direction: row-reverse; }
@keyframes mi { from{opacity:0;transform:translateY(4px)} to{opacity:1;transform:translateY(0)} }
.av {
  width: 30px; height: 30px; min-width: 30px; border-radius: 7px;
  display: flex; align-items: center; justify-content: center;
  font-size: 0.68rem; font-weight: 600; flex-shrink: 0;
}
.av-u { background: #6366f1; color: #fff; }
.av-b { background: #18181b; border: 1px solid #27272a; color: #71717a; font-size: 13px; }
.mb   { max-width: 76%; display: flex; flex-direction: column; gap: 4px; }
.msg.u .mb { align-items: flex-end; }
.bub {
  padding: 0.62rem 0.88rem; border-radius: 11px;
  font-size: 0.865rem; line-height: 1.65; word-break: break-word;
}
.bu { background: #6366f1; color: #fff; border-bottom-right-radius: 3px; }
.bb { background: #18181b; border: 1px solid #27272a; color: #e4e4e7; border-bottom-left-radius: 3px; }
.src { display: flex; flex-wrap: wrap; gap: 3px; }
.chip {
  display: inline-flex; align-items: center; gap: 3px;
  font-size: 0.65rem; font-weight: 500; color: #6366f1;
  background: rgba(99,102,241,.1); border: 1px solid rgba(99,102,241,.18);
  border-radius: 4px; padding: 2px 6px;
}

/* ── Output card ── */
.ocard {
  background: #18181b; border: 1px solid #27272a; border-radius: 10px;
  padding: 1.3rem 1.6rem; margin-top: 0.6rem;
  font-size: 0.875rem; line-height: 1.75; color: #e4e4e7;
}
.ocard h2 { color: #fafafa; font-weight: 600; font-size: 0.97rem; margin: 1em 0 0.35em; }
.ocard h3 { color: #a1a1aa; font-weight: 500; font-size: 0.88rem; margin: 0.9em 0 0.3em; }
.ocard ul { padding-left: 1.1em; }
.ocard li { margin-bottom: 3px; }
.ocard strong { color: #fafafa; }
.ocard code { background:#27272a; border-radius:4px; padding:1px 5px; font-size:0.82em; color:#a5b4fc; }

/* ── Empty state ── */
.empty {
  display:flex; flex-direction:column; align-items:center;
  justify-content:center; gap:8px; padding:3rem 2rem;
  color:#3f3f46; text-align:center;
}
.et { font-size:0.87rem; font-weight:500; color:#52525b; }
.es { font-size:0.75rem; color:#3f3f46; line-height:1.6; }

/* ── Download button ── */
.stDownloadButton > button {
  background: transparent !important; border: 1px solid #27272a !important;
  color: #a1a1aa !important; font-size: 0.77rem !important;
  padding: 0.38rem 0.75rem !important; width: auto !important; border-radius: 6px !important;
}
.stDownloadButton > button:hover { background: #18181b !important; color: #d4d4d8 !important; }

/* ── Divider ── */
hr { border-color: #1f1f23 !important; margin: 0.85rem 0 !important; }
.stSpinner > div { border-top-color: #6366f1 !important; }
</style>
""", unsafe_allow_html=True)


# ── Text extraction ─────────────────────────────────────────────────────────────
def get_icon(filename: str) -> tuple:
    ext = pathlib.Path(filename).suffix.lower()
    return {
        ".pdf":  ("📕", "t-pdf"),
        ".docx": ("📘", "t-docx"), ".doc": ("📘", "t-docx"),
        ".pptx": ("📙", "t-pptx"), ".ppt": ("📙", "t-pptx"),
        ".txt":  ("📄", "t-txt"),
    }.get(ext, ("📄", "t-txt"))


def extract_pages(file_bytes: bytes, filename: str) -> list:
    ext = pathlib.Path(filename).suffix.lower()
    if ext == ".pdf":
        reader, out = PdfReader(io.BytesIO(file_bytes)), []
        for i, page in enumerate(reader.pages, 1):
            t = (page.extract_text() or "").strip()
            if t:
                out.append({"text": t, "page": i})
        return out
    if ext == ".docx" and HAS_DOCX:
        try:
            doc  = DocxDocument(io.BytesIO(file_bytes))
            full = "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
            return [{"text": full[i:i+CHUNK_SIZE*2], "page": n+1}
                    for n, i in enumerate(range(0, max(len(full),1), CHUNK_SIZE*2))]
        except Exception as e:
            print(f"Error reading docx: {e}")
            return []
    if ext == ".pptx" and HAS_PPTX:
        try:
            prs, out = Presentation(io.BytesIO(file_bytes)), []
            for i, slide in enumerate(prs.slides, 1):
                texts = [s.text.strip() for s in slide.shapes if hasattr(s,"text") and s.text.strip()]
                if texts:
                    out.append({"text": "\n".join(texts), "page": i})
            return out
        except Exception as e:
            print(f"Error reading pptx: {e}")
            return []
    if ext == ".ppt" and HAS_PPT2TXT:
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".ppt") as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            
            parsed = ppt2txt.process(tmp_path)
            os.remove(tmp_path)

            if "content" in parsed:
                full = "\n".join(parsed["content"].values())
                return [{"text": full[i:i+CHUNK_SIZE*2], "page": n+1}
                        for n, i in enumerate(range(0, max(len(full),1), CHUNK_SIZE*2))]
        except Exception as e:
            print(f"Error reading ppt: {e}")
            return []
    if ext == ".txt":
        full = file_bytes.decode("utf-8", errors="ignore")
        return [{"text": full[i:i+CHUNK_SIZE*2], "page": n+1}
                for n, i in enumerate(range(0, max(len(full),1), CHUNK_SIZE*2))]
    return []


# ── ChromaDB ───────────────────────────────────────────────────────────────────
@st.cache_resource(show_spinner=False)
def get_collection():
    client = chromadb.PersistentClient(path="chroma_db")
    ef = embedding_functions.SentenceTransformerEmbeddingFunction(model_name="all-MiniLM-L6-v2")
    return client.get_or_create_collection(name="course_material", embedding_function=ef)

collection = get_collection()


def ingest(file_bytes: bytes, filename: str) -> int:
    data     = collection.get()
    existing = {m["source"] for m in data["metadatas"]} if data["metadatas"] else set()
    if filename in existing:
        return 0
    chunks, base = [], collection.count()
    for pg in extract_pages(file_bytes, filename):
        text, page = pg["text"], pg["page"]
        start = 0
        while start < len(text):
            c = text[start:start+CHUNK_SIZE].strip()
            if c:
                chunks.append({"text": c, "source": filename, "page": page})
            start += CHUNK_SIZE - OVERLAP
    if not chunks:
        return 0
    collection.add(
        ids=[f"doc_{base+i}" for i in range(len(chunks))],
        documents=[c["text"] for c in chunks],
        metadatas=[{"source": c["source"], "page": c["page"]} for c in chunks],
    )
    return len(chunks)


def indexed_docs() -> list:
    if collection.count() == 0:
        return []
    data = collection.get()
    if not data["metadatas"]:
        return []
    seen, names = set(), []
    for m in data["metadatas"]:
        if m["source"] not in seen:
            seen.add(m["source"])
            names.append(m["source"])
    return sorted(names)


# ── LLM ────────────────────────────────────────────────────────────────────────
def call_llm(prompt: str, provider: str, model: str) -> str:
    if provider == "Groq":
        from groq import Groq
        r = Groq(api_key=GROQ_API_KEY).chat.completions.create(
            model=model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3,
        )
        return r.choices[0].message.content
    else:
        import google.generativeai as genai
        genai.configure(api_key=GEMINI_API_KEY)
        return genai.GenerativeModel(model).generate_content(prompt).text


def retrieve(query: str, k: int = 8):
    if collection.count() == 0:
        return []
    r = collection.query(query_texts=[query], n_results=min(k, collection.count()))
    return list(zip(r["documents"][0], r["metadatas"][0], r["distances"][0]))


def ask(query: str, provider: str, model: str):
    hits = [h for h in retrieve(query, 8) if h[2] < DISTANCE_THRESHOLD]
    if not hits:
        return "This topic doesn't appear in your uploaded course material.", []
    ctx = "\n\n".join(f"[{m['source']}, p.{m['page']}]\n{doc}" for doc,m,_ in hits)
    prompt = (
        "You are a precise study assistant. Answer ONLY from the context below.\n"
        "If not found, say: \"This topic isn't covered in your course material.\"\n"
        "Be thorough and exam-relevant. Use markdown. End with Sources listing files and pages.\n\n"
        f"Context:\n{ctx}\n\nQuestion: {query}\nAnswer:"
    )
    return call_llm(prompt, provider, model), hits


def gen_notes(topic: str, provider: str, model: str) -> str:
    hits = retrieve(topic, 15)
    if not hits:
        return "No content found for this topic."
    ctx = "\n\n".join(f"[{m['source']} p.{m['page']}]\n{doc}" for doc,m,_ in hits)
    return call_llm(
        f'Write structured exam-revision notes on "{topic}" using ONLY the context below.\n'
        "Use ## headings, bullet points, **bold** key terms, highlight high-yield facts. Be concise.\n\n"
        f"Context:\n{ctx}\n\nRevision Notes:",
        provider, model
    )


def gen_qa(topic: str, n: int, provider: str, model: str) -> str:
    hits = retrieve(topic, 15)
    if not hits:
        return "No content found for this topic."
    ctx = "\n\n".join(f"[{m['source']} p.{m['page']}]\n{doc}" for doc,m,_ in hits)
    return call_llm(
        f"Generate {n} exam questions with model answers on \"{topic}\" using ONLY the context.\n"
        "Format: **Q[n].** [question]\n**Answer:** [answer] *(Source: page X)*\n"
        f"Mix short-answer and MCQ.\n\nContext:\n{ctx}\n\nQuestions and Answers:",
        provider, model
    )


# ── Session state ───────────────────────────────────────────────────────────────
for k, v in [("history",[]),("notes_out",""),("notes_topic",""),("qa_out",""),("qa_topic","")]:
    if k not in st.session_state:
        st.session_state[k] = v


# ════════════════════════════════════════════════════════════════════════════════
#  TOP NAV
# ════════════════════════════════════════════════════════════════════════════════
docs      = indexed_docs()
has_docs  = bool(docs)
dot_color = "#22c55e" if has_docs else "#3f3f46"
doc_label = f"{len(docs)} file{'s' if len(docs)!=1 else ''} indexed" if has_docs else "No files indexed"

st.markdown(f"""
<div class="top-nav">
  <div class="nav-brand">
    <div class="nav-logo">📖</div>
    <span class="nav-name">StudyBot</span>
    <span class="nav-tag">RAG</span>
  </div>
  <span class="status-pill">
    <span class="s-dot" style="background:{dot_color};"></span>
    {doc_label}
  </span>
</div>
""", unsafe_allow_html=True)


# ════════════════════════════════════════════════════════════════════════════════
#  TWO-COLUMN LAYOUT  — widgets placed DIRECTLY in columns, no HTML wrappers
# ════════════════════════════════════════════════════════════════════════════════
left, right = st.columns([1, 2.8])


# ╔══════════════════╗
# ║   LEFT PANEL     ║  ← real Streamlit widgets, no wrapping divs
# ╚══════════════════╝
with left:

    # ── Brand mark ──────────────────────────────────────────────────────────────
    st.markdown("""
    <div style="margin-bottom:1rem;">
      <span style="font-size:0.67rem;font-weight:600;color:#52525b;
                   text-transform:uppercase;letter-spacing:.08em;">AI Provider</span>
    </div>
    """, unsafe_allow_html=True)

    # ── Provider selector (direct Streamlit widget) ──────────────────────────────
    provider_options = []
    if GROQ_API_KEY:   provider_options.append("Groq")
    if GEMINI_API_KEY: provider_options.append("Gemini")
    if not provider_options: provider_options = ["Groq", "Gemini"]

    provider    = st.selectbox("Provider", provider_options, key="prov_sel", label_visibility="collapsed")
    model_map   = GROQ_MODELS if provider == "Groq" else GEMINI_MODELS
    model_label = st.selectbox("Model", list(model_map.keys()), key="model_sel", label_visibility="collapsed")
    model_id    = model_map[model_label]
    # Save to session_state so they're accessible from any column / tab
    st.session_state["_provider"] = provider
    st.session_state["_model"]    = model_id

    if provider == "Groq":
        st.markdown('<div class="prov-tag p-groq">⚡ Groq &mdash; Free &amp; fast</div>', unsafe_allow_html=True)
        if not GROQ_API_KEY:
            st.markdown('<div class="banner bw">Add GROQ_API_KEY to .env</div>', unsafe_allow_html=True)
    else:
        st.markdown('<div class="prov-tag p-gemini">🔵 Google Gemini</div>', unsafe_allow_html=True)
        if not GEMINI_API_KEY:
            st.markdown('<div class="banner bw">Add GEMINI_API_KEY to .env</div>', unsafe_allow_html=True)

    st.divider()

    # ── Upload label + format badges (HTML only — no widget) ────────────────────
    st.markdown("""
    <span class="sec-label">Upload Documents</span>
    <div class="fmt-row">
      <span class="fmt-badge">PDF</span>
      <span class="fmt-badge">DOCX</span>
      <span class="fmt-badge">PPTX</span>
      <span class="fmt-badge">TXT</span>
    </div>
    """, unsafe_allow_html=True)

    # ── File uploader (DIRECT Streamlit widget — no div wrapper) ────────────────
    uploaded = st.file_uploader(
        "Upload your course files",
        type=["pdf", "docx", "pptx", "ppt", "txt"],
        accept_multiple_files=True,
        label_visibility="collapsed",
    )

    # Process uploads
    if uploaded:
        added = 0
        with st.spinner("Indexing files…"):
            for f in uploaded:
                raw = f.read()
                n   = ingest(raw, f.name)
                if n > 0:
                    added += 1
                    dest = PDF_DIR / f.name
                    if not dest.exists():
                        dest.write_bytes(raw)
        if added:
            st.markdown(
                f'<div class="banner bo">✓ Indexed {added} new file{"s" if added>1 else ""}.</div>',
                unsafe_allow_html=True,
            )
            st.rerun()
        else:
            st.markdown('<div class="banner bi">All files already indexed.</div>', unsafe_allow_html=True)

    st.divider()

    # ── Document list (HTML display only) ───────────────────────────────────────
    docs = indexed_docs()
    if docs:
        st.markdown('<span class="sec-label">Indexed Files</span>', unsafe_allow_html=True)
        for name in docs:
            icon, cls = get_icon(name)
            short = name if len(name) <= 24 else name[:21] + "…"
            st.markdown(f"""
            <div class="doc-item">
              <div class="doc-ico {cls}">{icon}</div>
              <span class="doc-name" title="{name}">{short}</span>
              <span class="doc-ok"></span>
            </div>""", unsafe_allow_html=True)

        st.markdown(f"""
        <div class="stats-grid">
          <div class="stat-card">
            <div class="stat-n">{len(docs)}</div>
            <div class="stat-l">Files</div>
          </div>
          <div class="stat-card">
            <div class="stat-n">{collection.count():,}</div>
            <div class="stat-l">Chunks</div>
          </div>
        </div>""", unsafe_allow_html=True)

        st.markdown("<div style='height:0.6rem;'></div>", unsafe_allow_html=True)

    else:
        st.markdown("""
        <div class="empty" style="padding:1.2rem 0.5rem;">
          <svg width="26" height="26" viewBox="0 0 24 24" fill="none" stroke="#3f3f46" stroke-width="1.5">
            <path d="M13 2H6a2 2 0 0 0-2 2v16a2 2 0 0 0 2 2h12a2 2 0 0 0 2-2V9z"/>
            <polyline points="13 2 13 9 20 9"/>
          </svg>
          <div class="et">No files yet</div>
          <div class="es">Upload above to begin</div>
        </div>""", unsafe_allow_html=True)

    # ── Clear button (DIRECT Streamlit widget) ───────────────────────────────────
    if st.button("Clear conversation", use_container_width=True, key="clear_btn"):
        st.session_state.history   = []
        st.session_state.notes_out = ""
        st.session_state.qa_out    = ""
        st.rerun()


# ╔══════════════════╗
# ║   RIGHT PANEL    ║
# ╚══════════════════╝
with right:

    if not has_docs:
        st.markdown('<div class="banner bw">Upload your course files in the left panel to begin.</div>', unsafe_allow_html=True)

    tab_chat, tab_notes, tab_exam = st.tabs(["Chat", "Revision Notes", "Exam Practice"])

    # ── CHAT ────────────────────────────────────────────────────────────────────
    with tab_chat:
        if st.session_state.history:
            for msg in st.session_state.history:
                if msg["role"] == "user":
                    st.markdown(f"""
                    <div class="msg u">
                      <div class="av av-u">You</div>
                      <div class="mb"><div class="bub bu">{msg["content"]}</div></div>
                    </div>""", unsafe_allow_html=True)
                else:
                    chips = ""
                    if msg.get("sources"):
                        chips = '<div class="src">' + "".join(
                            f'<span class="chip">📄 {s["source"]}, p.{s["page"]}</span>'
                            for s in msg["sources"]
                        ) + '</div>'
                    content = msg["content"].replace("\n", "<br>")
                    st.markdown(f"""
                    <div class="msg">
                      <div class="av av-b">AI</div>
                      <div class="mb">
                        <div class="bub bb">{content}</div>
                        {chips}
                      </div>
                    </div>""", unsafe_allow_html=True)
        else:
            st.markdown("""
            <div class="empty" style="min-height:260px;">
              <svg width="34" height="34" viewBox="0 0 24 24" fill="none" stroke="#3f3f46" stroke-width="1.5">
                <path d="M21 15a2 2 0 0 1-2 2H7l-4 4V5a2 2 0 0 1 2-2h14a2 2 0 0 1 2 2z"/>
              </svg>
              <div class="et">No messages yet</div>
              <div class="es">Answers come strictly from your uploaded documents.</div>
            </div>""", unsafe_allow_html=True)

        # Input row — direct widgets
        ic1, ic2 = st.columns([6, 1])
        with ic1:
            question = st.text_input("Your question",
                                     placeholder="Ask something from your course material…",
                                     key="chat_q")
        with ic2:
            st.markdown("<div style='padding-top:1.6rem;'>", unsafe_allow_html=True)
            send = st.button("Send", key="send_btn", use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)

        if send and question.strip():
            if not has_docs:
                st.markdown('<div class="banner bw">Upload documents in the left panel first.</div>', unsafe_allow_html=True)
            elif provider == "Groq" and not GROQ_API_KEY:
                st.markdown('<div class="banner be">GROQ_API_KEY missing in .env</div>', unsafe_allow_html=True)
            elif provider == "Gemini" and not GEMINI_API_KEY:
                st.markdown('<div class="banner be">GEMINI_API_KEY missing in .env</div>', unsafe_allow_html=True)
            else:
                st.session_state.history.append({"role": "user", "content": question})
                with st.spinner("Searching your documents…"):
                    try:
                        answer, hits = ask(question, provider, model_id)
                        sources = [{"source": m["source"], "page": m["page"]} for _,m,_ in hits]
                        st.session_state.history.append({"role":"assistant","content":answer,"sources":sources})
                    except Exception as e:
                        err = "Rate limit hit. Switch provider or wait." if ("quota" in str(e).lower() or "429" in str(e)) else str(e)
                        st.session_state.history.append({"role":"assistant","content":f"⚠️ {err}","sources":[]})
                st.rerun()

    # ── REVISION NOTES ───────────────────────────────────────────────────────────
    with tab_notes:
        st.markdown('<p style="font-size:0.8rem;color:#71717a;margin-bottom:0.75rem;">Enter a topic to generate structured exam-ready notes from your documents.</p>', unsafe_allow_html=True)

        nt = st.text_input("Topic for notes", placeholder="e.g. Cardiovascular Pharmacology, Chapter 4…", key="notes_in")
        nb = st.button("Generate Revision Notes", key="notes_btn", use_container_width=True)

        if nb:
            if not has_docs:
                st.markdown('<div class="banner bw">Upload documents first.</div>', unsafe_allow_html=True)
            elif not nt.strip():
                st.markdown('<div class="banner bw">Enter a topic.</div>', unsafe_allow_html=True)
            else:
                _provider = st.session_state.get("_provider", provider)
                _model    = st.session_state.get("_model",    model_id)
                with st.spinner(f'Generating notes on "{nt}"...'):
                    try:
                        st.session_state.notes_out   = gen_notes(nt, _provider, _model)
                        st.session_state.notes_topic = nt
                    except Exception as e:
                        st.session_state.notes_out = f"⚠️ Error: {e}"

        if st.session_state.notes_out:
            st.markdown('<div class="ocard">', unsafe_allow_html=True)
            st.markdown(st.session_state.notes_out)
            st.markdown('</div>', unsafe_allow_html=True)
            st.download_button("⬇ Download notes (.md)",
                               data=st.session_state.notes_out,
                               file_name=f"notes_{st.session_state.notes_topic}.md",
                               mime="text/markdown")
        else:
            st.markdown("""
            <div class="empty" style="min-height:200px;">
              <svg width="30" height="30" viewBox="0 0 24 24" fill="none" stroke="#3f3f46" stroke-width="1.5">
                <path d="M14 2H6a2 2 0 0 0-2 2v16a2 2 0 0 0 2 2h12a2 2 0 0 0 2-2V8z"/>
                <polyline points="14 2 14 8 20 8"/>
                <line x1="16" y1="13" x2="8" y2="13"/><line x1="16" y1="17" x2="8" y2="17"/>
              </svg>
              <div class="et">No notes generated</div>
              <div class="es">Type a topic above and click Generate</div>
            </div>""", unsafe_allow_html=True)

    # ── EXAM PRACTICE ────────────────────────────────────────────────────────────
    with tab_exam:
        st.markdown('<p style="font-size:0.8rem;color:#71717a;margin-bottom:0.75rem;">Generate practice exam questions with model answers from your documents.</p>', unsafe_allow_html=True)

        et    = st.text_input("Topic for exam questions", placeholder="e.g. Renal Physiology, Antibiotics…", key="exam_in")
        num_q = st.slider("Number of questions", 3, 20, 8, key="num_q")
        eb    = st.button("Generate Exam Questions", key="exam_btn", use_container_width=True)

        if eb:
            if not has_docs:
                st.markdown('<div class="banner bw">Upload documents first.</div>', unsafe_allow_html=True)
            elif not et.strip():
                st.markdown('<div class="banner bw">Enter a topic.</div>', unsafe_allow_html=True)
            else:
                _provider = st.session_state.get("_provider", provider)
                _model    = st.session_state.get("_model",    model_id)
                with st.spinner(f'Generating {num_q} questions on "{et}"...'):
                    try:
                        st.session_state.qa_out   = gen_qa(et, num_q, _provider, _model)
                        st.session_state.qa_topic = et
                    except Exception as e:
                        st.session_state.qa_out = f"⚠️ Error: {e}"

        if st.session_state.qa_out:
            st.markdown('<div class="ocard">', unsafe_allow_html=True)
            st.markdown(st.session_state.qa_out)
            st.markdown('</div>', unsafe_allow_html=True)
            st.download_button("⬇ Download Q&A (.md)",
                               data=st.session_state.qa_out,
                               file_name=f"exam_qa_{st.session_state.qa_topic}.md",
                               mime="text/markdown")
        else:
            st.markdown("""
            <div class="empty" style="min-height:200px;">
              <svg width="30" height="30" viewBox="0 0 24 24" fill="none" stroke="#3f3f46" stroke-width="1.5">
                <circle cx="12" cy="12" r="10"/>
                <path d="M9.09 9a3 3 0 0 1 5.83 1c0 2-3 3-3 3"/>
                <line x1="12" y1="17" x2="12.01" y2="17"/>
              </svg>
              <div class="et">No questions generated</div>
              <div class="es">Type a topic above and click Generate</div>
            </div>""", unsafe_allow_html=True)
